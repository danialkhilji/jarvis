from pathlib import Path
import config

SKILL_NAME = "edit_powerpoint"
SKILL_DESCRIPTION = (
    "Read or edit PowerPoint files (.pptx). "
    "Actions: read, add_slide, replace_text. "
    "Returns slide text or confirmation of changes."
)


def run(
    file_path: str,
    action: str,
    output_path: str = "",
    slide_title: str = "",
    slide_body: str = "",
    search_text: str = "",
    replace_text: str = "",
    slide_number: int = 0,
) -> str:
    """
    Work with PowerPoint files.

    Args:
        file_path: Absolute path to the .pptx file.
        action: One of 'read', 'add_slide', 'replace_text'.
        output_path: For write actions; defaults to overwriting original.
        slide_title: Title for new slide (for 'add_slide').
        slide_body: Body text for new slide (for 'add_slide').
        search_text: Text to find (for 'replace_text').
        replace_text: Replacement text (for 'replace_text').
        slide_number: 1-based slide index to read (0 = all slides, for 'read').
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx is not installed. Run: pip install python-pptx"

    src = Path(file_path)
    if not src.exists():
        return f"Error: File not found: {file_path}"
    if src.suffix.lower() != ".pptx":
        return "Error: Only .pptx files are supported."

    allowed = any(str(src).startswith(d) for d in config.ALLOWED_DIRECTORIES)
    if not allowed:
        return "Error: Access denied. Path must be under an allowed directory."

    prs = Presentation(str(src))

    if action == "read":
        slides = prs.slides
        target = [slides[slide_number - 1]] if slide_number else slides
        lines = []
        for i, slide in enumerate(target, start=slide_number or 1):
            texts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
            lines.append(f"Slide {i}: {' | '.join(texts)}")
        return "\n".join(lines) if lines else "No text found in slides."

    dest = Path(output_path) if output_path else src

    if action == "add_slide":
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        if slide_title and slide.shapes.title:
            slide.shapes.title.text = slide_title
        if slide_body and len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_body

    elif action == "replace_text":
        if not search_text:
            return "Error: 'replace_text' requires search_text."
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if search_text in run.text:
                                run.text = run.text.replace(search_text, replace_text)
                                count += 1
        if count == 0:
            return f"Text '{search_text}' not found in presentation."

    else:
        return f"Error: Unknown action '{action}'. Choose from: read, add_slide, replace_text"

    prs.save(str(dest))
    return f"PowerPoint saved: {dest}"
